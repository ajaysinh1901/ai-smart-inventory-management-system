# -*- coding: utf-8 -*-
"""
Builds the MCA Semester-II viva PowerPoint for the
"AI Smart Inventory Management System" (MERN stack) project.

White + blue theme. Follows the college example template's 14-slide order:
Title, Introduction, Problem Statement, Project Objectives, Technologies Used,
System Architecture, React JS, Node.js & Express.js, MongoDB, API Testing,
Advantages, Future Scope, Conclusion, Thank You.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------- theme
NAVY    = RGBColor(0x0B, 0x2A, 0x4A)   # deep navy headings
PRIMARY = RGBColor(0x1D, 0x4E, 0xD8)   # primary blue
BLUE    = RGBColor(0x25, 0x63, 0xEB)   # bright blue
SKY     = RGBColor(0x0E, 0xA5, 0xE9)   # sky accent
TEXT    = RGBColor(0x1F, 0x29, 0x37)   # body text
MUTED   = RGBColor(0x64, 0x74, 0x8B)   # muted gray
CARDBG  = RGBColor(0xEF, 0xF4, 0xFF)   # light blue card
CARDBD  = RGBColor(0xC7, 0xD9, 0xFB)   # card border
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BANDBG  = RGBColor(0xF3, 0xF7, 0xFF)   # faint blue panel

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

PROJECT = "AI Smart Inventory Management System"


# ----------------------------------------------------------------------------- helpers
def add_slide():
    return prs.slides.add_slide(BLANK)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    fill(sp, color)
    sp.shadow.inherit = False
    return sp


def card(slide, x, y, w, h, fillc=CARDBG, border=CARDBD):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fillc
    sp.line.color.rgb = border
    sp.line.width = Pt(1)
    sp.shadow.inherit = False
    # softer corner
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    return sp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def para(tf, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    return p


def header(slide, kicker, title):
    """Standard content-slide header: top blue strip, kicker, title, underline."""
    # top accent strip
    rect(slide, 0, 0, SW, Inches(0.16), PRIMARY)
    # kicker
    _, tf = textbox(slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.35))
    set_run(para(tf, True).add_run(), kicker.upper(), 13, SKY, bold=True)
    tf.paragraphs[0].runs[0].font.name = FONT
    # title
    _, tf = textbox(slide, Inches(0.7), Inches(0.82), Inches(11.9), Inches(0.8))
    set_run(para(tf, True).add_run(), title, 32, NAVY, bold=True)
    # underline bar
    bar = rect(slide, Inches(0.72), Inches(1.62), Inches(1.1), Inches(0.07), BLUE)
    return Inches(1.95)   # content top


def footer(slide, n):
    _, tf = textbox(slide, Inches(0.7), Inches(7.05), Inches(9), Inches(0.3))
    set_run(para(tf, True).add_run(), PROJECT, 10.5, MUTED)
    _, tf = textbox(slide, Inches(12.0), Inches(7.05), Inches(0.7), Inches(0.3))
    p = para(tf, True); p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), str(n), 10.5, MUTED, bold=True)


def bullets(slide, items, top, left=Inches(0.85), width=Inches(11.6),
            size=19, gap=12, marker=PRIMARY):
    """items: list of (text, bold_text_or_None) or plain strings.
       Supports inline bold via list of (segment, is_bold) tuples."""
    _, tf = textbox(slide, left, top, width, Inches(4.8))
    for i, item in enumerate(items):
        p = para(tf, i == 0)
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        # marker
        r = p.add_run()
        set_run(r, "▪  ", 19, marker, bold=True)
        if isinstance(item, str):
            set_run(p.add_run(), item, size, TEXT)
        else:
            for seg, isb in item:
                set_run(p.add_run(), seg, size, NAVY if isb else TEXT, bold=isb)
    return tf


def card_grid(slide, cards, top, cols=3, gap=Inches(0.35),
              left=Inches(0.7), right=Inches(0.7), row_h=Inches(1.55), rows=None):
    avail = SW - left - right
    cw = (avail - gap * (cols - 1)) / cols
    n = len(cards)
    rows = rows or ((n + cols - 1) // cols)
    for idx, (htitle, body) in enumerate(cards):
        r, c = divmod(idx, cols)
        x = left + c * (cw + gap)
        y = top + r * (row_h + Inches(0.28))
        card(slide, x, y, cw, row_h)
        _, tf = textbox(slide, x + Inches(0.22), y + Inches(0.18),
                        cw - Inches(0.44), row_h - Inches(0.36))
        set_run(para(tf, True).add_run(), htitle, 16, PRIMARY, bold=True)
        p = para(tf); p.space_before = Pt(4); p.line_spacing = 1.08
        set_run(p.add_run(), body, 13, TEXT)


# ============================================================================= SLIDE 1 — TITLE
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
# blue bands top & bottom
rect(s, 0, 0, SW, Inches(0.55), PRIMARY)
rect(s, 0, SH - Inches(0.35), SW, Inches(0.35), NAVY)
# thin sky line under top band
rect(s, 0, Inches(0.55), SW, Inches(0.06), SKY)

# logo placeholder (rounded blue square with monogram)
logo = card(s, Inches(5.96), Inches(1.05), Inches(1.4), Inches(1.4), PRIMARY, PRIMARY)
_, tf = textbox(s, Inches(5.96), Inches(1.05), Inches(1.4), Inches(1.4))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "AI", 40, WHITE, bold=True)

# title
_, tf = textbox(s, Inches(1.0), Inches(2.75), Inches(11.33), Inches(1.5))
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "AI Smart Inventory Management System", 40, NAVY, bold=True)
p = para(tf); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(4)
set_run(p.add_run(), "Using the MERN Stack", 26, PRIMARY, bold=True)

# tagline
_, tf = textbox(s, Inches(1.5), Inches(4.35), Inches(10.33), Inches(0.6))
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(),
        "An intelligent platform for inventory, GST billing & AI-driven business insights",
        16, MUTED, italic=True)

# meta block
_, tf = textbox(s, Inches(1.5), Inches(5.15), Inches(10.33), Inches(1.6))
for i, (a, b) in enumerate([
    ("MCA Semester-II Project", None),
    ("R.B. Institute of Management Studies (RBIMS), Ahmedabad", None),
    ("Affiliated to Gujarat Technological University (GTU)", None),
    ("Guide:  Prof. Himani Khodifad", None),
    ("Presented by:  ____________________     Enrolment No:  ____________", None),
    ("Academic Year:  2025 - 26", None),
]):
    p = para(tf, i == 0); p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(3)
    set_run(p.add_run(), a, 14.5 if i == 0 else 13,
            NAVY if i == 0 else TEXT, bold=(i == 0))


# ============================================================================= SLIDE 2 — INTRODUCTION
s = add_slide()
top = header(s, "Introduction", "Introduction")
bullets(s, [
    [("The ", False), ("AI Smart Inventory Management System", True),
     (" is a full-stack web application built on the MERN stack to digitise the daily operations of small retail businesses.", False)],
    [("It brings ", False), ("inventory, GST billing, customer credit (khata) and business analytics", True),
     (" together on a single, easy-to-use platform.", False)],
    [("The system adds ", False), ("Artificial Intelligence", True),
     (" using the Google Gemini API for demand forecasting and smart reorder suggestions.", False)],
    [("It is designed specifically for ", False), ("Indian kirana stores, traders and wholesalers", True),
     (" who currently rely on paper registers and khata books.", False)],
    [("The goal is to replace manual guesswork with ", False), ("real-time, data-driven decisions", True),
     (" through a simple interface for non-technical users.", False)],
], top)
footer(s, 2)


# ============================================================================= SLIDE 3 — PROBLEM STATEMENT
s = add_slide()
top = header(s, "Problem Statement", "Problem Statement")
bullets(s, [
    [("Millions of small shops still manage stock using ", False),
     ("paper registers or basic spreadsheets", True), (", which are slow and error-prone.", False)],
    [("Manual billing causes ", False), ("GST calculation errors", True),
     (", lost invoices and compliance difficulties.", False)],
    [("Customer credit (“udhaar”) is tracked in ", False), ("physical khata books", True),
     (" — entries get lost and disputes arise.", False)],
    [("Shop owners have ", False), ("no clear visibility", True),
     (" into fast-selling products, what to reorder, or stock lost to shrinkage.", False)],
    [("Existing software is often ", False), ("too expensive, too complex, or not built", True),
     (" for the Indian small-business workflow.", False)],
], top)
footer(s, 3)


# ============================================================================= SLIDE 4 — PROJECT OBJECTIVES
s = add_slide()
top = header(s, "Goals", "Project Objectives")
card_grid(s, [
    ("Digitise Inventory", "Replace manual registers with real-time stock tracking, units of measure, low-stock alerts and shrinkage records."),
    ("Simplify GST Billing", "One-tap Quick Sale billing with automatic GST computation, PDF invoices and UPI payment QR codes."),
    ("Add Intelligence", "Use Google Gemini AI to forecast demand, suggest reorders and surface actionable business insights."),
    ("Digital Khata", "Track customer credit and repayments digitally — a ledger that never loses an entry."),
    ("Analytics & Reports", "Dashboards, sales trends, scheduled reports and Tally-compatible export for accountants."),
    ("Accessible to All", "Multi-language interface, guided onboarding and a simple UI for non-technical shop owners."),
], top, cols=3, row_h=Inches(1.95))
footer(s, 4)


# ============================================================================= SLIDE 5 — TECHNOLOGIES USED
s = add_slide()
top = header(s, "Technology", "Technologies Used")
rows = [
    ("MongoDB", "NoSQL document database for flexible data storage"),
    ("Express.js", "Backend web framework for building the REST API"),
    ("React.js", "Frontend library for a fast, responsive single-page UI"),
    ("Node.js", "JavaScript runtime environment that powers the server"),
    ("Google Gemini API", "AI for demand forecasting and smart insights"),
    ("Tesseract.js (OCR)", "Scans supplier bills to auto-extract products"),
    ("JWT, bcrypt & Zod", "Authentication, password hashing and validation"),
    ("Thunder Client", "API testing of GET, POST, PUT and DELETE requests"),
]
# two-column table-like layout using cards
left = Inches(0.7); colw = Inches(5.9); gap = Inches(0.35); rowh = Inches(0.92)
for i, (tech, desc) in enumerate(rows):
    r, c = divmod(i, 2)
    x = left + c * (colw + gap)
    y = top + r * (rowh + Inches(0.18))
    card(s, x, y, colw, rowh)
    # tech name bullet
    rect(s, x + Inches(0.2), y + Inches(0.33), Inches(0.13), Inches(0.13), BLUE)
    _, tf = textbox(s, x + Inches(0.5), y + Inches(0.13), colw - Inches(0.7), rowh - Inches(0.26))
    set_run(para(tf, True).add_run(), tech, 16, NAVY, bold=True)
    p = para(tf); p.space_before = Pt(1)
    set_run(p.add_run(), desc, 12.5, TEXT)
footer(s, 5)


# ============================================================================= SLIDE 6 — SYSTEM ARCHITECTURE
s = add_slide()
top = header(s, "Design", "System Architecture")
intro = textbox(s, Inches(0.85), top, Inches(11.6), Inches(0.5))[1]
set_run(para(intro, True).add_run(),
        "A three-tier architecture with clear separation between presentation, logic and data.",
        15, MUTED, italic=True)

tiers = [
    ("Frontend Layer", "React 19 + Vite + Tailwind",
     ["Dashboard, Inventory, Quick Sale", "Khata, AI Insights, Analytics", "Recharts charts + i18n languages"]),
    ("Backend Layer", "Node.js + Express 5",
     ["REST API (controllers → services)", "JWT auth, Zod validation", "AI, OCR, PDF & cron jobs"]),
    ("Database Layer", "MongoDB + Mongoose",
     ["Products & Stock Adjustments", "Sales, Customers & Khata", "Suppliers, Alerts & Settings"]),
]
ty = top + Inches(0.65)
tw = Inches(3.65); th = Inches(3.45); tgap = Inches(0.55)
tleft = Inches(0.7)
for i, (name, tech, items) in enumerate(tiers):
    x = tleft + i * (tw + tgap)
    card(s, x, ty, tw, th, BANDBG, CARDBD)
    # header band
    rect(s, x, ty, tw, Inches(0.7), PRIMARY, MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = textbox(s, x + Inches(0.2), ty + Inches(0.1), tw - Inches(0.4), Inches(0.55))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = para(tf, True); p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), name, 17, WHITE, bold=True)
    _, tf = textbox(s, x + Inches(0.2), ty + Inches(0.85), tw - Inches(0.4), Inches(0.4))
    p = para(tf, True); p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), tech, 13, SKY, bold=True)
    _, tf = textbox(s, x + Inches(0.3), ty + Inches(1.35), tw - Inches(0.6), Inches(2.0))
    for j, it in enumerate(items):
        p = para(tf, j == 0); p.space_after = Pt(8); p.line_spacing = 1.05
        set_run(p.add_run(), "▪ ", 13, BLUE, bold=True)
        set_run(p.add_run(), it, 13, TEXT)
    # arrow between tiers
    if i < 2:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                x + tw + Inches(0.07), ty + th/2 - Inches(0.18),
                                Inches(0.42), Inches(0.36))
        fill(ar, SKY); ar.shadow.inherit = False

# bottom note
_, tf = textbox(s, Inches(0.7), ty + th + Inches(0.12), Inches(11.9), Inches(0.4))
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Client ↔ Server communication over REST APIs (Axios / JSON)  •  Mongoose ODM ↔ MongoDB",
        13, MUTED, bold=True)
footer(s, 6)


# ============================================================================= concept slide helper (React / Node / Mongo)
def concept_slide(n, kicker, title, lead, points, role_title, role_body):
    s = add_slide()
    top = header(s, kicker, title)
    # lead line
    _, tf = textbox(s, Inches(0.85), top, Inches(11.6), Inches(0.55))
    set_run(para(tf, True).add_run(), lead, 16, PRIMARY, bold=True, italic=True)
    # left bullets
    bullets(s, points, top + Inches(0.6), left=Inches(0.85), width=Inches(7.2), size=16, gap=10)
    # right role card
    cx = Inches(8.45); cw = Inches(4.1)
    card(s, cx, top + Inches(0.55), cw, Inches(3.6))
    _, tf = textbox(s, cx + Inches(0.3), top + Inches(0.8), cw - Inches(0.6), Inches(3.2))
    set_run(para(tf, True).add_run(), role_title, 16, PRIMARY, bold=True)
    p = para(tf); p.space_before = Pt(8); p.line_spacing = 1.2
    set_run(p.add_run(), role_body, 14, TEXT)
    footer(s, n)
    return s


# ============================================================================= SLIDE 7 — REACT JS
concept_slide(
    7, "Frontend Library", "React JS",
    "The 'V' (View) of the application — what the shop owner sees and clicks.",
    [
        [("A ", False), ("component-based", True), (" JavaScript library for building user interfaces.", False)],
        [("Uses a ", False), ("virtual DOM", True), (" for fast, efficient screen updates.", False)],
        [("Built with ", False), ("React 19 + Vite", True), (" for instant load and hot-reload during development.", False)],
        [("State managed with the ", False), ("Context API", True), (" (Auth, Theme, Language) — no Redux needed.", False)],
        [("Styled using ", False), ("Tailwind CSS", True), (" with charts drawn by Recharts.", False)],
    ],
    "Role in this project",
    "React renders every page — Dashboard, Inventory, Quick Sale, Khata, AI Insights and Analytics. "
    "Reusable UI components (cards, modals, KPI strips) keep the interface consistent, while i18next "
    "provides multi-language support for non-English users.",
)


# ============================================================================= SLIDE 8 — NODE & EXPRESS
concept_slide(
    8, "Backend", "Node.js and Express.js",
    "The server side — the brain that processes every request.",
    [
        [("Node.js", True), (" is a JavaScript runtime that lets us run JS on the server.", False)],
        [("It is ", False), ("event-driven and non-blocking", True), (", ideal for handling many requests.", False)],
        [("Express.js", True), (" is a minimal web framework built on top of Node.js.", False)],
        [("It exposes a ", False), ("REST API", True), (" using a clean controllers → services pattern.", False)],
        [("Handles ", False), ("auth, validation, AI, OCR, PDF invoices and cron jobs", True), (".", False)],
    ],
    "Role in this project",
    "Express powers 17 modular route groups under /api/v1. Each request flows through middleware "
    "(JWT auth, Zod validation, rate limiting) into controllers and service functions that contain the "
    "business logic — GST math, khata balances and AI calls — before talking to MongoDB.",
)


# ============================================================================= SLIDE 9 — MONGODB
concept_slide(
    9, "Database", "MongoDB",
    "The data layer — where everything is stored.",
    [
        [("A ", False), ("NoSQL", True), (" database that stores data as flexible JSON-like documents.", False)],
        [("Stores all ", False), ("user and application data", True), (" for the system.", False)],
        [("Flexible ", False), ("document-based storage", True), (" adapts easily as the app grows.", False)],
        [("Connected through ", False), ("Mongoose ODM", True), (" for schemas and validation.", False)],
        [("Money & weight use ", False), ("Decimal128", True), (" to avoid floating-point rounding errors.", False)],
    ],
    "Role in this project",
    "MongoDB holds 12+ collections — Products, Sales, Customers, KhataEntries, Suppliers, "
    "StockAdjustments, Alerts and Settings. Mongoose models enforce structure and keep currency precise, "
    "making the data reliable enough for real billing, not just a demo.",
)


# ============================================================================= SLIDE 10 — API TESTING
s = add_slide()
top = header(s, "Quality", "API Testing")
bullets(s, [
    [("Thunder Client", True), (" (a VS Code extension) is used to test the backend REST API.", False)],
    [("Every endpoint is tested for ", False), ("GET, POST, PUT and DELETE", True), (" requests.", False)],
    [("Responses are verified for correct ", False), ("status codes, data shape and GST math", True), (".", False)],
    [("Both ", False), ("success and error cases", True), (" are checked — e.g. invalid input is rejected.", False)],
    [("Automated ", False), ("Node smoke tests", True), (" further verify endpoints during development.", False)],
], top, size=18)
# method pills
methods = [("GET", "Read data"), ("POST", "Create"), ("PUT", "Update"), ("DELETE", "Remove")]
px = Inches(0.85); py = Inches(5.1); pw = Inches(2.85); gap = Inches(0.25)
for i, (m, d) in enumerate(methods):
    x = px + i * (pw + gap)
    card(s, x, py, pw, Inches(1.1), CARDBG)
    _, tf = textbox(s, x, py + Inches(0.16), pw, Inches(0.8))
    p = para(tf, True); p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), m, 20, PRIMARY, bold=True)
    p = para(tf); p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), d, 13, MUTED)
footer(s, 10)


# ============================================================================= SLIDE 11 — ADVANTAGES
s = add_slide()
top = header(s, "Benefits", "Advantages")
card_grid(s, [
    ("All-in-One Platform", "Inventory, billing, khata and analytics in a single application — no juggling registers."),
    ("Saves Time", "One-tap Quick Sale billing and OCR bill scanning make daily work much faster."),
    ("Accurate GST & Money", "Automatic GST and Decimal128 math remove calculation and rounding errors."),
    ("Smarter Decisions", "AI demand forecasting and smart alerts tell owners what to reorder and when."),
    ("No Lost Credit", "Digital khata keeps an exact record of who owes what — no disputes."),
    ("Simple & Accessible", "Clean UI, guided onboarding and multi-language support for non-technical users."),
    ("Secure", "JWT authentication, hashed passwords, input validation and rate limiting."),
    ("Affordable & Scalable", "Open MERN stack keeps it low-cost and easy to extend for more shops."),
], top, cols=4, row_h=Inches(2.1))
footer(s, 11)


# ============================================================================= SLIDE 12 — FUTURE SCOPE
s = add_slide()
top = header(s, "Roadmap", "Future Scope")
bullets(s, [
    [("Mobile app", True), (" (React Native / PWA) so billing works at the counter on any phone.", False)],
    [("WhatsApp / SMS integration", True), (" for sending invoices and automatic khata payment reminders.", False)],
    [("Barcode scanning", True), (" at the billing counter for packaged goods.", False)],
    [("Multi-store / multi-user workspaces", True), (" with role-based permissions for staff.", False)],
    [("Offline-first mode", True), (" with sync — billing must never stop when the internet does.", False)],
    [("Deeper AI", True), (" — price optimisation, festival-season demand models and theft anomaly detection.", False)],
], top, size=18)
footer(s, 12)


# ============================================================================= SLIDE 13 — CONCLUSION
s = add_slide()
top = header(s, "Summary", "Conclusion")
bullets(s, [
    [("The system successfully ", False), ("digitises the complete daily workflow", True),
     (" of a small retail business — inventory, billing, credit and analysis — in one app.", False)],
    [("AI integration is ", False), ("practical, not decorative", True),
     (": forecasts and alerts directly answer “what should I reorder and when?”", False)],
    [("The ", False), ("MERN architecture", True),
     (" with a clean controllers → services → models structure is modular, testable and easy to extend.", False)],
    [("Careful engineering — ", False), ("Decimal128 money math, Zod validation, JWT security, cron automation", True),
     (" — makes it reliable for real shops.", False)],
    [("The project demonstrates ", False), ("end-to-end full-stack skills", True),
     (": database design, REST APIs, modern React and AI/OCR integration.", False)],
], top, size=18)
footer(s, 13)


# ============================================================================= SLIDE 14 — THANK YOU
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, Inches(0.55), PRIMARY)
rect(s, 0, Inches(0.55), SW, Inches(0.06), SKY)
rect(s, 0, SH - Inches(0.35), SW, Inches(0.35), NAVY)

logo = card(s, Inches(6.06), Inches(2.05), Inches(1.2), Inches(1.2), PRIMARY, PRIMARY)
_, tf = textbox(s, Inches(6.06), Inches(2.05), Inches(1.2), Inches(1.2))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "AI", 34, WHITE, bold=True)

_, tf = textbox(s, Inches(1.5), Inches(3.5), Inches(10.33), Inches(1.2))
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Thank You", 58, NAVY, bold=True)
p = para(tf); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
set_run(p.add_run(), "Questions & Discussion", 20, PRIMARY)

_, tf = textbox(s, Inches(1.5), Inches(5.7), Inches(10.33), Inches(0.5))
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), PROJECT + "  •  MERN Stack Project  •  2025 - 26", 14, MUTED)


# ----------------------------------------------------------------------------- save
out = "presentation/AI-Smart-Inventory-Presentation.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
